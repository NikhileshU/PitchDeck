"""Peer agreement — card-schema.md §5.1.

The contract says both renderers stack content from the top of the content
area using one shared height model (validate._block_h), so a block lands at
the same top edge in both. Nothing enforced that until this file: §5.1 was a
document, and the 0.00pt result was a one-off checkpoint script.

These tests are the falsifiable version. If a renderer starts placing blocks
by its own arithmetic, this goes red.
"""

import pytest
from pptx import Presentation

import render_pptx  # noqa: E402
from validate import _block_h, _text_lines  # noqa: E402

from conftest import theme_for

EMU_PER_PT = 12700
TOLERANCE_PT = 0.01
CANVAS_W = 960.0


def expected_tops(ir, theme):
    """Recompute every content block's top edge from the shared model alone —
    deliberately NOT by calling into either renderer."""
    pad = theme["space"]["cardPad"]
    gap = theme["space"]["blockGap"]
    sc, lh = theme["type"]["scale"], theme["type"]["lineHeight"]
    width = CANVAS_W - 2 * pad
    out = {}
    for card in ir["cards"]:
        if card.get("layout") != "content":
            continue
        title = card.get("title") or ""
        y = pad
        if title:
            y += _text_lines(title, sc["cardTitle"], width) * sc["cardTitle"] * lh + gap
        rows = []
        for b in card.get("blocks") or []:
            rows.append((str(b["id"]), round(y, 4)))
            y += _block_h(b, theme, width, 1.0) + gap
        out[card["id"]] = rows
    return out


def actual_tops(pptx_path):
    """Shape names are `<blockid>:<type>`; that naming is what makes the
    contract checkable, including the invisible marker _b_columns emits."""
    prs = Presentation(str(pptx_path))
    out = []
    for slide in prs.slides:
        tops = {}
        for sh in slide.shapes:
            name = sh.name or ""
            if ":" in name:
                tops.setdefault(name.split(":")[0], sh.top / EMU_PER_PT)
        out.append(tops)
    return out


class TestPeerGeometry:
    def test_every_block_is_named(self, renderable_ir, tmp_path):
        """A block the PPTX doesn't name is a block the peer check cannot see.
        Regression guard for R9-H1: `columns` was silently unnamed, so the one
        card with nested geometry was the one the checkpoint skipped."""
        path, ir = renderable_ir
        theme = theme_for(ir["meta"]["theme"])
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, theme, str(out))

        slides = actual_tops(out)
        expected = expected_tops(ir, theme)
        for card, slide_tops in zip(ir["cards"], slides):
            if card.get("layout") != "content":
                continue
            for bid, _ in expected[card["id"]]:
                assert bid in slide_tops, (
                    f"{path.stem} {card['id']}: block {bid} has no named shape"
                )

    def test_top_edges_match_shared_model(self, renderable_ir, tmp_path):
        """§5.1 clause 3: PPTX places frames at the model's exact pt positions."""
        path, ir = renderable_ir
        theme = theme_for(ir["meta"]["theme"])
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, theme, str(out))

        slides = actual_tops(out)
        expected = expected_tops(ir, theme)
        for card, slide_tops in zip(ir["cards"], slides):
            if card.get("layout") != "content":
                continue
            for bid, want in expected[card["id"]]:
                got = slide_tops[bid]
                assert abs(got - want) < TOLERANCE_PT, (
                    f"{path.stem} {card['id']}/{bid}: expected {want}pt, got {got}pt"
                )

    def test_block_order_preserved(self, renderable_ir, tmp_path):
        """Clause 3 again: HTML reaches the geometry by flow and may wrap
        differently, but neither renderer may REORDER blocks."""
        path, ir = renderable_ir
        theme = theme_for(ir["meta"]["theme"])
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, theme, str(out))

        slides = actual_tops(out)
        expected = expected_tops(ir, theme)
        for card, slide_tops in zip(ir["cards"], slides):
            if card.get("layout") != "content":
                continue
            ir_order = [bid for bid, _ in expected[card["id"]]]
            drawn = sorted((b for b in ir_order), key=lambda b: slide_tops[b])
            assert drawn == ir_order, (
                f"{path.stem} {card['id']}: PPTX order {drawn} != IR order {ir_order}"
            )

    def test_slide_count_equals_card_count(self, renderable_ir, tmp_path):
        path, ir = renderable_ir
        theme = theme_for(ir["meta"]["theme"])
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, theme, str(out))
        assert len(Presentation(str(out)).slides) == len(ir["cards"])


class TestPeerGeometryAcrossThemes:
    """Type scale and spacing are identical across themes on purpose
    (R11-N4), so the same IR must land at the same coordinates in all three.
    If someone gives warm a roomier cardPad, this is what catches it."""

    def test_geometry_is_theme_invariant(self, all_blocks, tmp_path):
        baselines = {}
        for name in ("slate", "warm", "mono"):
            theme = theme_for(name)
            out = tmp_path / f"{name}.pptx"
            render_pptx.render(all_blocks, theme, str(out))
            baselines[name] = expected_tops(all_blocks, theme)
        assert baselines["slate"] == baselines["warm"] == baselines["mono"]
