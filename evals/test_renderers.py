"""Renderer invariants — properties, not golden bytes.

Deliberately NOT byte-comparison. HTML output is ~13KB of markup and PPTX is a
zip of XML; pinning either means rebaselining on every cosmetic change and
learning nothing. These assert the properties the spec actually promises, so
they survive refactors and still catch the failures that shipped bugs before.
"""

import re
import zipfile

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

import render_html
import render_pptx


class TestHtmlInvariants:
    """css="" throughout: these assert on markup, not on layout, so rendering
    unstyled is the deliberate choice. Omitting css entirely now raises (R13-M2)
    — one card per page is a property of base.css, not of the renderer."""

    def test_renders_every_fixture_every_theme(self, renderable_ir, any_theme, tmp_path):
        path, ir = renderable_ir
        name, theme = any_theme
        out = tmp_path / "deck.html"
        render_html.render(ir, theme, str(out), css="")
        assert out.stat().st_size > 0

    def test_no_javascript(self, renderable_ir, slate, tmp_path):
        """Invariant: no JS in rendered HTML. A theme value that escaped the
        <style> block would land here as an injected <script>."""
        _, ir = renderable_ir
        out = tmp_path / "deck.html"
        render_html.render(ir, slate, str(out), css="")
        html = out.read_text(encoding="utf-8")
        assert "<script" not in html.lower()
        assert "onclick" not in html.lower()

    def test_self_contained_no_local_paths(self, renderable_ir, slate, tmp_path):
        """'Self-contained' means the file survives being emailed. R1-H6 was
        exactly this: srcs rewritten to absolute local paths."""
        _, ir = renderable_ir
        out = tmp_path / "deck.html"
        render_html.render(ir, slate, str(out), css="")
        html = out.read_text(encoding="utf-8")
        assert "file://" not in html
        for src in re.findall(r'<img[^>]+src="([^"]+)"', html):
            assert src.startswith("data:"), f"non-embedded image src: {src[:60]}"

    def test_card_count_matches_ir(self, renderable_ir, slate, tmp_path):
        _, ir = renderable_ir
        out = tmp_path / "deck.html"
        render_html.render(ir, slate, str(out), css="")
        html = out.read_text(encoding="utf-8")
        assert html.count('class="card ') == len(ir["cards"])


class TestPptxInvariants:
    @staticmethod
    def _render(ir, theme, tmp_path):
        # ir_dir: fixtures carry relative image srcs, as deck.json does on disk
        from conftest import FIXTURES
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, theme, str(out), ir_dir=FIXTURES)
        return out

    def test_clause_6_no_autofit_anywhere(self, renderable_ir, slate, tmp_path):
        """card-schema.md §5.1 clause 6. PowerPoint's shrink-to-fit would absorb
        overflow that §5 requires to surface as a Tier-1 error, so a deck that
        validates clean must never render at silently reduced type.

        Covers table cells and decorative autoshapes too — the first phase-9
        run failed precisely because python-pptx gives autoshapes an implicit
        MIDDLE-anchored text frame."""
        _, ir = renderable_ir
        prs = Presentation(str(self._render(ir, slate, tmp_path)))
        frames = 0
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.has_text_frame:
                    frames += 1
                    assert sh.text_frame.auto_size == MSO_AUTO_SIZE.NONE, (
                        f"{sh.name}: auto_size is {sh.text_frame.auto_size}"
                    )
                if sh.has_table:
                    for row in sh.table.rows:
                        for cell in row.cells:
                            frames += 1
                            assert cell.text_frame.auto_size == MSO_AUTO_SIZE.NONE
        assert frames > 0

    def test_clause_6_written_into_xml(self, renderable_ir, slate, tmp_path):
        """Stronger than the API check: assert the flag is actually serialised.
        `auto_size = None` would remove the element and mean 'inherit', which is
        the trap clause 6 exists to close."""
        _, ir = renderable_ir
        path = self._render(ir, slate, tmp_path)
        with zipfile.ZipFile(path) as z:
            slides = [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml", n)]
            assert slides
            for n in slides:
                xml = z.read(n).decode("utf-8")
                assert "normAutofit" not in xml, f"{n} carries normAutofit"
                assert "spAutoFit" not in xml, f"{n} carries spAutoFit"
                for anchor in re.findall(r'anchor="(\w+)"', xml):
                    assert anchor == "t", f"{n} anchors {anchor}, not top"

    def test_charts_are_native_not_pictures(self, all_blocks, slate, tmp_path):
        """No screenshotting, ever. A rasterised chart would have to exist as a
        file in ppt/media/, so media count must equal the number of genuine
        `image` blocks — anything above that is a chart or table drawn as a
        picture."""
        path = self._render(all_blocks, slate, tmp_path)
        from validate import _iter_blocks
        want = sum(
            1
            for card in all_blocks["cards"]
            for b, _ in _iter_blocks(card.get("blocks"))
            if b.get("type") == "image"
        )
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            media = [n for n in names if n.startswith("ppt/media/")]
            charts = [n for n in names if n.startswith("ppt/charts/chart")]
        assert charts, "no native chart parts — charts may have been rasterised"
        # <= not ==: python-pptx dedupes identical images across slides, so
        # two blocks pointing at one file yield one media part.
        assert len(media) <= want, (
            f"{len(media)} media files for {want} image blocks — "
            "something non-image was rasterised"
        )

    def test_text_is_selectable(self, all_blocks, slate, tmp_path):
        path = self._render(all_blocks, slate, tmp_path)
        with zipfile.ZipFile(path) as z:
            runs = sum(
                z.read(n).decode("utf-8").count("<a:t>")
                for n in z.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml", n)
            )
        assert runs > 20, f"only {runs} text runs — text may be an image"

    def test_canvas_is_16x9_at_960x540pt(self, all_blocks, slate, tmp_path):
        prs = Presentation(str(self._render(all_blocks, slate, tmp_path)))
        assert prs.slide_width / 12700 == pytest.approx(960.0)
        assert prs.slide_height / 12700 == pytest.approx(540.0)

    def test_theme_font_family_reaches_powerpoint(self, all_blocks, any_theme, tmp_path):
        """A theme that renders in HTML but loses its family in PPTX is a peer
        divergence users see immediately."""
        name, theme = any_theme
        if name == "bad-contrast":
            pytest.skip("not a shipped theme")
        prs = Presentation(str(self._render(all_blocks, theme, tmp_path)))
        expected = theme["type"]["family"].split(",")[0].strip().strip("'\"")
        seen = {
            r.font.name
            for slide in prs.slides
            for sh in slide.shapes
            if sh.has_text_frame
            for p in sh.text_frame.paragraphs
            for r in p.runs
            if r.font.name
        }
        assert seen == {expected}, f"{name}: expected {expected}, got {seen}"


class TestLineBudget:
    """Both renderers are capped at 400 *code* lines (BUILD-SPEC §15, R13-N7).

    The cap exists to stop the block switch accumulating logic that belongs
    elsewhere. Counting raw lines measured the wrong thing: it taxed comments and
    docstrings equally, and a bug fix was once paid for by deleting explanation.
    Blanks, comments and docstrings are free; statements are not."""

    CAP = 400

    @staticmethod
    def _code_lines(path):
        import ast
        src = path.read_text(encoding="utf-8")
        lines = src.splitlines()
        skip = sum(1 for line in lines if not line.strip() or line.strip().startswith("#"))
        for node in ast.walk(ast.parse(src)):
            if isinstance(node, (ast.Module, ast.FunctionDef, ast.ClassDef)):
                if ast.get_docstring(node):
                    doc = node.body[0]
                    skip += doc.end_lineno - doc.lineno + 1
        return len(lines) - skip

    @pytest.mark.parametrize("name", ["render_html.py", "render_pptx.py"])
    def test_renderer_is_within_the_code_budget(self, name):
        from conftest import ROOT
        n = self._code_lines(ROOT / "scripts" / name)
        assert n <= self.CAP, (
            f"{name} is {n} code lines, over the {self.CAP} cap — if the block "
            "switch has grown logic, move it, do not delete comments to fit"
        )


class TestLibraryImageResolution:
    """R13-M1. Image srcs resolve relative to deck.json (block-types.md), and that
    step used to live in each CLI's main(), so `render()` called as a library got
    FileNotFoundError on a relative src — or, if the caller resolved paths itself,
    an HTML file carrying absolute local paths, which is the R1-H6 defect wearing
    a different hat. Both renderers now take `ir_dir` as a parameter, the same fix
    that closed R4-H4 in validate.py. The IR here is loaded straight from disk with
    its relative srcs untouched, because that is what a library caller passes."""

    IMAGE_FIXTURE = "all-blocks.json"

    def _raw_ir(self):
        from conftest import FIXTURES, load
        return load(FIXTURES / self.IMAGE_FIXTURE), FIXTURES

    def test_html_embeds_relative_srcs_given_ir_dir(self, slate, tmp_path):
        ir, base = self._raw_ir()
        out = tmp_path / "deck.html"
        render_html.render(ir, slate, str(out), css="", ir_dir=base)
        html = out.read_text(encoding="utf-8")
        srcs = re.findall(r'<img[^>]+src="([^"]+)"', html)
        assert srcs, "fixture no longer has image blocks — pick another"
        for s in srcs:
            assert s.startswith("data:"), f"not embedded: {s[:60]}"

    def test_pptx_resolves_relative_srcs_given_ir_dir(self, slate, tmp_path):
        ir, base = self._raw_ir()
        out = tmp_path / "deck.pptx"
        render_pptx.render(ir, slate, str(out), ir_dir=base)
        assert Presentation(str(out)).slides

    @pytest.mark.parametrize("renderer,suffix", [(render_html, "html"), (render_pptx, "pptx")])
    def test_without_ir_dir_a_relative_src_fails_loudly(self, renderer, suffix, slate, tmp_path):
        """Omitting ir_dir is legitimate — it means 'srcs are already absolute'.
        What it must never do is silently emit a deck with a broken image."""
        ir = {
            "schema": "1.0",
            "meta": {"title": "t", "archetype": "product-demo",
                     "theme": "slate", "aspect": "16:9"},
            "cards": [{"id": "c1", "layout": "content", "title": "Revenue grew 12% in Q3",
                       "blocks": [{"id": "b1", "type": "image",
                                   "src": "nowhere/missing.png", "alt": "x"}],
                       "notes": "n"}],
        }
        with pytest.raises((FileNotFoundError, ValueError)):
            renderer.render(ir, slate, str(tmp_path / f"deck.{suffix}"))


class TestRendererErrorParity:
    """Both renderers must reject the same bad IR with the same-shaped error.
    Geometry parity is worth little if the peers disagree on what is invalid."""

    BAD = [
        ("empty series", {"id": "b1", "type": "chart", "chart": "bar",
                          "data": {"categories": ["a"], "series": []}}),
        ("non-numeric values", {"id": "b1", "type": "chart", "chart": "bar",
                                "data": {"categories": ["a", "b"],
                                         "series": [{"name": "x", "values": [1, None]}]}}),
        ("unknown block type", {"id": "b1", "type": "timeline", "text": "nope"}),
        ("nested columns", {"id": "b1", "type": "columns", "children": [
            [{"id": "b2", "type": "columns", "children": [
                [{"id": "b4", "type": "text", "text": "a"}],
                [{"id": "b5", "type": "text", "text": "b"}]]}],
            [{"id": "b3", "type": "text", "text": "c"}]]}),
    ]

    @pytest.mark.parametrize("label,block", BAD, ids=[b[0] for b in BAD])
    def test_both_renderers_reject(self, label, block, slate, tmp_path):
        ir = {
            "schema": "1.0",
            "meta": {"title": "t", "archetype": "product-demo",
                     "theme": "slate", "aspect": "16:9"},
            "cards": [{"id": "c1", "layout": "content",
                       "title": "Revenue grew 12% in Q3",
                       "blocks": [block], "notes": "n"}],
        }
        with pytest.raises(ValueError) as html_err:
            render_html.render(ir, slate, str(tmp_path / "d.html"), css="")
        with pytest.raises(ValueError) as pptx_err:
            render_pptx.render(ir, slate, str(tmp_path / "d.pptx"))
        assert "c1" in str(html_err.value) and "c1" in str(pptx_err.value)
