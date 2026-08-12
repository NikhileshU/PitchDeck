#!/usr/bin/env python3
"""render_pptx.py — peer renderer: deck.json IR -> native, editable PPTX.

Pure: render(ir, theme, out_path). Geometry follows card-schema.md §5.1 — blocks
in IR order at the shared stacking model's pt positions (heights imported from
validate._block_h, the single source). Clause 6: every text frame gets explicit
no-autofit (MSO_AUTO_SIZE.NONE -> <a:noAutofit/>) and MSO_ANCHOR.TOP — PowerPoint
must never shrink text to absorb overflow. Charts and tables are native objects.
No raw OOXML. Relative image srcs resolve against ir_dir when given (block-types.md).
"""
import argparse, copy, json, sys
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from validate import _block_h, _text_lines

CHART_KIND = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
              "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}

def _rgb(h): return RGBColor.from_string(str(h).strip().lstrip("#"))

def _sty(theme):
    c, t, s = theme["color"], theme["type"], theme["space"]
    return {"bg": _rgb(c["bg"]), "surface": _rgb(c["surface"]), "text": _rgb(c["text"]),
            "muted": _rgb(c["muted"]), "accent": _rgb(c["accent"]),
            "series": [_rgb(x) for x in c["series"]],
            "tone": {k: _rgb(v) for k, v in c["tone"].items()},
            "family": str(t["family"]).split(",")[0].strip(),
            "sc": t["scale"], "lh": t["lineHeight"], "pad": s["cardPad"], "gap": s["blockGap"]}

def _setup_tf(tf):
    # §5.1 clause 6: explicit no-autofit + top anchor on every frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def _frame(shapes, x, y, w, h, name=None):
    tb = shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    if name:
        tb.name = name
    return tb, _setup_tf(tb.text_frame)

def _para(tf, text, size, color, st, bold=False, italic=False, align=None, after=None):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run()
    r.text = str(text)
    f = r.font
    f.size, f.name, f.bold, f.italic = Pt(size), st["family"], bold, italic
    f.color.rgb = color
    if align is not None:
        p.alignment = align
    if after is not None:
        p.space_after = Pt(after)
    return p

def _rect(shapes, shape, x, y, w, h, fill, name=None):
    sh = shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    _setup_tf(sh.text_frame)  # clause 6 applies to autoshapes' implicit frames too
    if name:
        sh.name = name
    return sh

# ---- blocks (exactly 10) ---------------------------------------------------

def _b_text(shapes, b, st, theme, x, y, w, h, name):
    _, tf = _frame(shapes, x, y, w, h, name)
    _para(tf, b["text"], st["sc"]["body"], st["text"], st, bold=bool(b.get("emphasis")))

def _b_bullets(shapes, b, st, theme, x, y, w, h, name):
    _, tf = _frame(shapes, x, y, w, h, name)
    for it in b["items"]:
        _para(tf, f"•  {it}", st["sc"]["body"], st["text"], st, after=st["gap"] / 2)

def _b_kpi(shapes, b, st, theme, x, y, w, h, name, scale=1.0):
    grp = shapes.add_group_shape()
    grp.name = name
    sc, lh = st["sc"], st["lh"]
    vh = sc["title"] * scale * lh
    _, tf = _frame(grp.shapes, x, y, w, vh)
    _para(tf, b["value"], sc["title"] * scale, st["accent"], st, bold=True)
    _, tf = _frame(grp.shapes, x, y + vh, w, sc["body"] * lh)
    _para(tf, b["label"], sc["body"], st["text"], st)
    if b.get("delta"):
        _, tf = _frame(grp.shapes, x, y + vh + sc["body"] * lh, w, sc["caption"] * lh)
        _para(tf, b["delta"], sc["caption"], st["muted"], st)

def _b_chart(shapes, b, st, theme, x, y, w, h, name):
    cats, series = b["data"].get("categories"), b["data"].get("series")
    if not cats or not series:  # parity with render_html: fail loudly, never a blank frame
        raise ValueError("chart needs at least one category and one series")
    data = CategoryChartData()
    data.categories = [str(c) for c in cats]
    for s in series:
        if not all(isinstance(v, (int, float)) and not isinstance(v, bool)
                   for v in s.get("values") or []):
            raise ValueError(f"series {s.get('name')!r} has non-numeric values")
        data.add_series(str(s["name"]), tuple(float(v) for v in s["values"]))
    # chart box = model height minus the caption strip — follows _block_h wherever it moves
    cap_h = st["sc"]["caption"] * st["lh"] + st["gap"] / 2 if b.get("caption") else 0
    gf = shapes.add_chart(CHART_KIND[b["chart"]], Pt(x), Pt(y), Pt(w), Pt(h - cap_h), data)
    gf.name = name
    ch = gf.chart
    ch.has_title = False
    nseries = len(series)
    if b["chart"] == "pie":
        for i, pt in enumerate(ch.plots[0].series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = st["series"][i % len(st["series"])]
    else:
        for i, s in enumerate(ch.series):
            col = st["series"][i % len(st["series"])]
            if b["chart"] == "line":
                s.format.line.color.rgb = col
            else:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = col
    ch.has_legend = b["chart"] == "pie" or nseries > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(st["sc"]["caption"])
        ch.legend.font.color.rgb = st["muted"]
    for axis in ("category_axis", "value_axis"):
        try:
            ax = getattr(ch, axis)
            ax.tick_labels.font.size = Pt(st["sc"]["caption"])
            ax.tick_labels.font.color.rgb = st["muted"]
            ax.has_major_gridlines = False
        except (ValueError, AttributeError):
            pass  # pie has no axes
    if b.get("caption"):
        _, tf = _frame(shapes, x, y + (h - cap_h) + st["gap"] / 2, w,
                       st["sc"]["caption"] * st["lh"])
        _para(tf, b["caption"], st["sc"]["caption"], st["muted"], st)

def _b_table(shapes, b, st, theme, x, y, w, h, name):
    headers, rows = b["headers"], b["rows"]
    gf = shapes.add_table(len(rows) + 1, len(headers), Pt(x), Pt(y), Pt(w), Pt(h))
    gf.name = name
    tbl = gf.table
    tbl.first_row = tbl.horz_banding = False
    for col in tbl.columns:
        col.width = Pt(w / len(headers))
    for ri, cells in enumerate([headers] + list(rows)):
        for ci, val in enumerate(cells):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = st["surface"] if ri == 0 else st["bg"]
            tf = _setup_tf(cell.text_frame)
            tf.margin_left = tf.margin_right = Pt(st["gap"] / 2)
            tf.margin_top = tf.margin_bottom = Pt(st["gap"] / 3)
            if ri == 0:
                _para(tf, val, st["sc"]["caption"], st["muted"], st, bold=True)
            else:
                _para(tf, val, st["sc"]["body"], st["text"], st)

def _b_callout(shapes, b, st, theme, x, y, w, h, name):
    box = _rect(shapes, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, st["surface"], name)
    _rect(shapes, MSO_SHAPE.RECTANGLE, x, y, 4.5, h, st["tone"][b["tone"]])
    tf = _setup_tf(box.text_frame)
    tf.margin_left = tf.margin_right = Pt(st["gap"])
    tf.margin_top = tf.margin_bottom = Pt(st["gap"] / 1.5)
    _para(tf, b["text"], st["sc"]["body"], st["text"], st)

def _b_quote(shapes, b, st, theme, x, y, w, h, name):
    _rect(shapes, MSO_SHAPE.RECTANGLE, x, y, 4.5, h, st["accent"])
    _, tf = _frame(shapes, x + st["gap"], y, w - st["gap"], h, name)
    _para(tf, b["text"], st["sc"]["body"], st["text"], st, italic=True)
    if b.get("attribution"):
        _para(tf, b["attribution"], st["sc"]["caption"], st["muted"], st, after=None)

def _b_image(shapes, b, st, theme, x, y, w, h, name):
    iw, ih = Image.open(b["src"]).size
    fr, imr = w / h, iw / ih
    if b.get("fit", "contain") == "cover":
        pic = shapes.add_picture(b["src"], Pt(x), Pt(y), Pt(w), Pt(h))
        if imr > fr:
            pic.crop_left = pic.crop_right = (1 - fr / imr) / 2
        else:
            pic.crop_top = pic.crop_bottom = (1 - imr / fr) / 2
    else:
        pw, ph = (w, w / imr) if imr > fr else (h * imr, h)
        pic = shapes.add_picture(b["src"], Pt(x + (w - pw) / 2), Pt(y + (h - ph) / 2),
                                 Pt(pw), Pt(ph))
    pic.name = name

def _b_columns(shapes, b, st, theme, x, y, w, h, name):
    cols = b["children"]
    if any(cb.get("type") == "columns" for col in cols for cb in col):
        raise ValueError("columns may not nest inside columns")
    # invisible geometry marker so peer-agreement checks can locate this block
    tb = shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(0.5))
    tb.name = name
    _setup_tf(tb.text_frame)
    cw = (w - (len(cols) - 1) * st["gap"]) / len(cols)
    for ci, col in enumerate(cols):
        cy = y
        for cb in col:
            hh = _block_h(cb, theme, cw, 1.0)
            _block(shapes, cb, st, theme, x + ci * (cw + st["gap"]), cy, cw, hh)
            cy += hh + st["gap"]

def _b_divider(shapes, b, st, theme, x, y, w, h, name):
    _rect(shapes, MSO_SHAPE.RECTANGLE, x, y, w, 1, st["muted"], name)

_BLOCKS = {"text": _b_text, "bullets": _b_bullets, "kpi": _b_kpi, "chart": _b_chart,
           "table": _b_table, "callout": _b_callout, "quote": _b_quote,
           "image": _b_image, "columns": _b_columns, "divider": _b_divider}

def _block(shapes, b, st, theme, x, y, w, h):
    t = b.get("type")
    if t not in _BLOCKS:
        raise ValueError(f"unknown block type {t!r} (block {b.get('id')!r})")
    try:
        _BLOCKS[t](shapes, b, st, theme, x, y, w, h, f"{b.get('id')}:{t}")
    except (KeyError, ValueError) as e:  # IR faults; renderer bugs propagate raw
        raise ValueError(f"block {b.get('id')!r} ({t}): {e}") from e

# ---- cards -----------------------------------------------------------------

def _title_h(card, st, width):
    if not card.get("title"):
        return 0
    return _text_lines(card["title"], st["sc"]["cardTitle"], width) * st["sc"]["cardTitle"] * st["lh"]

def _card(prs, card, theme, st):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lay = card["layout"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = st["surface"] if lay == "section" else st["bg"]
    pad, gap, sc, lh = st["pad"], st["gap"], st["sc"], st["lh"]
    W = 960 - 2 * pad
    blocks = card.get("blocks") or []
    try:
        if lay == "content":
            th = _title_h(card, st, W)
            if th:
                _, tf = _frame(slide.shapes, pad, pad, W, th, "card-title")
                _para(tf, card["title"], sc["cardTitle"], st["text"], st, bold=True)
            y = pad + (th + gap if th else 0)
            for b in blocks:
                h = _block_h(b, theme, W, 1.0)
                _block(slide.shapes, b, st, theme, pad, y, W, h)
                y += h + gap
        elif lay in ("title", "section"):
            tl = _text_lines(card.get("title", ""), sc["title"], W) * sc["title"] * lh
            sub = sc["body"] * lh if card.get("subtitle") else 0
            rule = gap + 4.5 if lay == "title" else 0
            top = (540 - (tl + (gap + sub if sub else 0) + rule)) / 2
            _, tf = _frame(slide.shapes, pad, top, W, tl, "card-title")
            _para(tf, card.get("title", ""), sc["title"], st["text"], st, bold=True,
                  align=PP_ALIGN.CENTER)
            if sub:
                _, tf = _frame(slide.shapes, pad, top + tl + gap, W, sub, "card-subtitle")
                _para(tf, card["subtitle"], sc["body"], st["muted"], st, align=PP_ALIGN.CENTER)
            if lay == "title":
                _rect(slide.shapes, MSO_SHAPE.RECTANGLE, 480 - pad, top + tl + gap + sub,
                      pad * 2, 4.5, st["accent"])
        else:  # hero — §5.1 clause 4: image full-bleed + overlay; else centered stack
            img = next((b for b in blocks if b.get("type") == "image"), None)
            kpi = next((b for b in blocks if b.get("type") == "kpi"), None)
            tl = _text_lines(card.get("title", ""), sc["title"], W) * sc["title"] * lh \
                if card.get("title") else 0
            if img:
                _b_image(slide.shapes, dict(img, fit="cover"), st, theme, 0, 0, 960, 540,
                         f"{img.get('id')}:image")
                if tl:
                    _, tf = _frame(slide.shapes, pad, (540 - tl) / 2, W, tl, "card-title")
                    _para(tf, card["title"], sc["title"], st["text"], st, bold=True,
                          align=PP_ALIGN.CENTER)
            else:
                kh = (sc["title"] * 2 + sc["body"] + (sc["caption"] if kpi and kpi.get("delta")
                      else 0)) * lh if kpi else 0
                top = (540 - (kh + (gap + tl if tl else 0))) / 2
                if kpi:
                    _b_kpi(slide.shapes, kpi, st, theme, pad, top, W,
                           kh, f"{kpi.get('id')}:kpi", scale=2.0)
                if tl:
                    _, tf = _frame(slide.shapes, pad, top + kh + (gap if kh else 0), W, tl,
                                   "card-title")
                    _para(tf, card["title"], sc["title"], st["text"], st, bold=True,
                          align=PP_ALIGN.CENTER)
    except ValueError as e:
        raise ValueError(f"card {card.get('id')!r}: {e}") from e
    if card.get("notes"):
        slide.notes_slide.notes_text_frame.text = str(card["notes"])

# ---- entry -----------------------------------------------------------------

def render(ir, theme, out_path, ir_dir=None):
    """ir_dir is the directory relative image srcs resolve against (block-types.md).
    Pass it and render() resolves them itself; omit it and every src must already be
    absolute. It is a parameter, never discovered — same fix as validate.py's ir_dir
    (R4-H4), so a library caller runs the same path the CLI does."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(960), Pt(540)
    st = _sty(theme)
    if ir_dir is not None:
        ir = copy.deepcopy(ir)  # resolution must not mutate the caller's IR
        for card in ir["cards"]:
            _resolve_images(card.get("blocks"), Path(ir_dir).resolve())
    for card in ir["cards"]:
        _card(prs, card, theme, st)
    prs.save(out_path)

def _resolve_images(blocks, base):
    for b in blocks or []:
        if b.get("type") == "image" and not Path(str(b.get("src", ""))).is_absolute():
            p = base / b["src"]
            if not p.exists():
                raise FileNotFoundError(f"image not found for block {b.get('id')!r}: {p}")
            b["src"] = str(p)
        elif b.get("type") == "columns":
            for col in b.get("children") or []:
                _resolve_images(col, base)

def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--ir", required=True)
    ap.add_argument("--theme", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args(argv)
    # One guard over the whole run — see render_html.main: a CLI fails with a
    # message, never a traceback (R13-L5, widened). OSError covers an unreadable
    # --ir and an unwritable --out; ValueError covers malformed JSON and this
    # module's own raises; KeyError/TypeError cover valid JSON that is not a deck.
    try:
        ir = json.loads(Path(args.ir).read_text(encoding="utf-8"))
        theme = json.loads(Path(args.theme).read_text(encoding="utf-8"))
        Path(args.out).parent.mkdir(parents=True, exist_ok=True)
        render(ir, theme, args.out, ir_dir=Path(args.ir).resolve().parent)
    except (OSError, ValueError, KeyError, TypeError) as e:
        print(f"render_pptx: {type(e).__name__}: {e}", file=sys.stderr)
        return 1
    return 0

if __name__ == "__main__":
    sys.exit(main())
